# result_exports.py

from django.contrib.auth.decorators import login_required
from django.http import HttpResponse

from openpyxl import Workbook

from version1.matches_app.utils.results_utils import get_results_context


@login_required
def results_export_excel(request, team):

    competition_id = request.GET.get("competition", "all")
    date_from = request.GET.get("date_from") or None
    date_to = request.GET.get("date_to") or None

    results = get_results_context(
        team_code=team,
        competition_id=competition_id,
        date_from=date_from,
        date_to=date_to,
    )

    matches = results["past_matches"]

    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "Results"

    worksheet.append([
        "Date",
        "Competition",
        "Home Team",
        "Home Goals",
        "Away Goals",
        "Away Team",
        "Venue",
    ])

    for match in matches:
        worksheet.append([
            match.date,
            match.competition.name if match.competition else "",
            match.home_team.name,
            match.our_team_goals
            if match.home_team.id == match.our_team_id
            else match.opponent_goals,
            match.opponent_goals
            if match.home_team.id == match.our_team_id
            else match.our_team_goals,
            match.away_team.name,
            match.venue.name if match.venue else "",
        ])

    response = HttpResponse(
        content_type=(
            "application/vnd.openxmlformats-officedocument."
            "spreadsheetml.sheet"
        )
    )

    response["Content-Disposition"] = (
        f'attachment; filename="{team}_results.xlsx"'
    )

    workbook.save(response)

    return response






from django.contrib.auth.decorators import login_required
from django.http import HttpResponse

from reportlab.lib import colors
from reportlab.lib.pagesizes import landscape, A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import (
    SimpleDocTemplate,
    Table,
    TableStyle,
    Paragraph,
)

from version1.matches_app.utils.results_utils import get_results_context


@login_required
def results_export_pdf(request, team):

    competition_id = request.GET.get("competition", "all")
    date_from = request.GET.get("date_from") or None
    date_to = request.GET.get("date_to") or None

    results = get_results_context(
        team_code=team,
        competition_id=competition_id,
        date_from=date_from,
        date_to=date_to,
    )

    matches = results["past_matches"]

    response = HttpResponse(
        content_type="application/pdf"
    )

    response["Content-Disposition"] = (
        f'attachment; filename="{team}_results.pdf"'
    )

    document = SimpleDocTemplate(
        response,
        pagesize=landscape(A4),
    )

    styles = getSampleStyleSheet()

    elements = []

    elements.append(
        Paragraph(
            f"{team} - Match Results",
            styles["Title"],
        )
    )

    data = [[
        "Date",
        "Competition",
        "Home Team",
        "Score",
        "Away Team",
        "Venue",
    ]]

    for match in matches:

        if match.home_team.id == match.our_team_id:
            home_goals = match.our_team_goals
            away_goals = match.opponent_goals
        else:
            home_goals = match.opponent_goals
            away_goals = match.our_team_goals

        data.append([
            match.date.strftime("%d %b %Y"),
            match.competition.name if match.competition else "",
            match.home_team.name,
            f"{home_goals} - {away_goals}",
            match.away_team.name,
            match.venue.name if match.venue else "",
        ])

    table = Table(data)

    table.setStyle(
        TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.grey),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.black),
            ("ALIGN", (3, 1), (3, -1), "CENTER"),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
            ("PADDING", (0, 0), (-1, -1), 6),
        ])
    )

    elements.append(table)

    document.build(elements)

    return response









from pptx import Presentation
from pptx.util import Inches, Pt

from django.contrib.auth.decorators import login_required
from django.http import HttpResponse

from version1.matches_app.utils.results_utils import get_results_context


@login_required
def results_export_ppt(request, team):

    competition_id = request.GET.get("competition", "all")
    date_from = request.GET.get("date_from") or None
    date_to = request.GET.get("date_to") or None

    results = get_results_context(
        team_code=team,
        competition_id=competition_id,
        date_from=date_from,
        date_to=date_to,
    )

    matches = results["past_matches"]

    presentation = Presentation()

    # Title slide
    slide = presentation.slides.add_slide(
        presentation.slide_layouts[0]
    )

    slide.shapes.title.text = f"{team} Match Results"

    subtitle = slide.placeholders[1]

    subtitle.text = (
        f"{date_from or 'All dates'} - "
        f"{date_to or 'Today'}"
    )

    # Results slide
    slide = presentation.slides.add_slide(
        presentation.slide_layouts[5]
    )

    slide.shapes.title.text = "Match Results"

    rows = len(matches) + 1
    cols = 5

    table = slide.shapes.add_table(
        rows,
        cols,
        Inches(0.5),
        Inches(1.5),
        Inches(12.3),
        Inches(5.2),
    ).table

    headers = [
        "Date",
        "Opponent",
        "Result",
        "Competition",
        "Venue",
    ]

    for column, header in enumerate(headers):
        table.cell(0, column).text = header

    for row, match in enumerate(matches, start=1):

        if match.home_team.id == match.our_team_id:
            opponent = match.away_team.name
            our_goals = match.our_team_goals
            opponent_goals = match.opponent_goals
        else:
            opponent = match.home_team.name
            our_goals = match.our_team_goals
            opponent_goals = match.opponent_goals

        if our_goals > opponent_goals:
            result = f"W {our_goals}-{opponent_goals}"
        elif our_goals < opponent_goals:
            result = f"L {our_goals}-{opponent_goals}"
        else:
            result = f"D {our_goals}-{opponent_goals}"

        values = [
            match.date.strftime("%d %b %Y"),
            opponent,
            result,
            match.competition.name if match.competition else "",
            match.venue.name if match.venue else "",
        ]

        for column, value in enumerate(values):
            table.cell(row, column).text = str(value)

    response = HttpResponse(
        content_type=(
            "application/vnd.openxmlformats-officedocument."
            "presentationml.presentation"
        )
    )

    response["Content-Disposition"] = (
        f'attachment; filename="{team}_results.pptx"'
    )

    presentation.save(response)

    return response