import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment
from django.http import HttpResponse
from django.shortcuts import get_object_or_404, render
from django.utils.dateparse import parse_date
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

from version1.teams_app.models import Team
from version1.reports_app.models.previous_models import Result, Medical, Transition, Scouting, Performance, IndividualActionPlan
from version1.reports_app.views.daily_report_views.statistics_view import get_statistics_report
from version1.players_app.models import Player


from reportlab.lib import colors
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import mm
from reportlab.platypus import ( SimpleDocTemplate, Paragraph,Spacer,Table,TableStyle,PageBreak,)

from reportlab.lib.pagesizes import A4, landscape
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def get_full_name(player):
    """Return the full player name as: First + Second + Surname"""
    if not player:
        return ""
    return f"{player.name} {player.second_name} {player.surname}".strip()

def download_technical_report_ppt(request, team_id):
    team = get_object_or_404(Team, id=team_id)
    start = parse_date(request.GET.get("start_date", ""))
    end = parse_date(request.GET.get("end_date", ""))
    season = request.GET.get("season", "")

    results = Result.objects.filter(our_team=team)
    if start: results = results.filter(date__gte=start)
    if end: results = results.filter(date__lte=end)

    stats = get_statistics_report(
        filter_type=request.GET.get("filter", "all"),
        team=team, start_date=start, end_date=end
    )

    prs = Presentation()

    # ---------- TITLE ----------
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "TECHNICAL REPORT"
    slide.placeholders[1].text = (
        f"{team.name}\nPeriod: {start or 'All'} to {end or 'All'}"
    )

    # ---------- RESULTS ----------
    result_rows = []
    for r in results:
        scorers = []
        for s in (r.goal_scorers or "").split(","):
            s = s.strip()
            if s:
                try:
                    name, minute = s.rsplit(" ", 1)
                    p = Player.objects.filter(name__iexact=name).first()
                    scorers.append(
                        f"{get_full_name(p)} {minute}" if p else s
                    )
                except ValueError:
                    scorers.append(s)

        result_rows.append([
            str(r.date), str(r.competition_type),
            r.home_team.name, f"{r.home_score}-{r.away_score}",
            r.away_team.name, str(r.result), ", ".join(scorers)
        ])

    result_headers = [
        "Date", "Competition", "Home", "Score",
        "Away", "Result", "Goal Scorers"
    ]

    # Maximum 15 results per slide
    for start_row in range(0, len(result_rows), 15):
        chunk = result_rows[start_row:start_row + 15]

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = (
            "MATCH RESULTS" if start_row == 0
            else "MATCH RESULTS - CONTINUED"
        )

        table = slide.shapes.add_table(
            len(chunk) + 1, 7,
            Inches(.3), Inches(1.2),
            Inches(12.7), Inches(5.5)
        ).table

        for j, h in enumerate(result_headers):
            table.cell(0, j).text = h

        for i, row in enumerate(chunk, 1):
            for j, value in enumerate(row):
                table.cell(i, j).text = value

    # ---------- STATISTICS ----------
    headers = [
        "Player", "Pos",
        "Game", "Apps", "Starts", "Sub In", "Sub Out",
        "Goals", "Assists", "Pre Assists",
        "Total", "U11", "U13", "U15", "U17",
        "U20", "First", "National"
    ]

    stat_rows = []
    for r in stats:
        stat_rows.append([
            get_full_name(r.get("player")),
            str(r.get("position", "")),

            # MATCH DATA
            str(r.get("game_minutes", 0)),
            str(r.get("appearances", 0)),
            str(r.get("starts", 0)),
            str(r.get("sub_in", 0)),
            str(r.get("sub_out", 0)),
            str(r.get("goals", 0)),
            str(r.get("assists", 0)),
            str(r.get("pre_assists", 0)),

            # TRAINING DATA
            str(r.get("training_total", 0)),
            str(r.get("training_u11", 0)),
            str(r.get("training_u13", 0)),
            str(r.get("training_u15", 0)),
            str(r.get("training_u17", 0)),
            str(r.get("training_u20", 0)),
            str(r.get("training_first", 0)),
            str(r.get("training_national", 0)),
        ])

    # Maximum 17 players per slide
    for start_row in range(0, len(stat_rows), 17):
        chunk = stat_rows[start_row:start_row + 17]

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = (
            "PLAYER STATISTICS" if start_row == 0
            else "PLAYER STATISTICS - CONTINUED"
        )

        table = slide.shapes.add_table(
            len(chunk) + 1, 18,
            Inches(.2), Inches(1.2),
            Inches(12.9), Inches(5.5)
        ).table

        for j, h in enumerate(headers):
            table.cell(0, j).text = h

        for i, row in enumerate(chunk, 1):
            for j, value in enumerate(row):
                table.cell(i, j).text = value

    # ---------- DOWNLOAD ----------
    response = HttpResponse(
        content_type=(
            "application/vnd.openxmlformats-"
            "officedocument.presentationml.presentation"
        )
    )

    response["Content-Disposition"] = (
        f'attachment; filename="Technical_Report_'
        f'{team.age_group.code}_{season or "ALL"}.pptx"'
    )

    prs.save(response)
    return response
