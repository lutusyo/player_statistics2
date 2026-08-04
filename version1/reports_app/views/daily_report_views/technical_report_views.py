import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment
from django.http import HttpResponse
from django.shortcuts import get_object_or_404, render
from django.utils.dateparse import parse_date
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

from version1.teams_app.models import Team
from version1.reports_app.models.previous_models import Result, Medical, Transition, Scouting, Performance, IndividualActionPlan
from version1.reports_app.views.daily_report_views.statistics_view import get_statistics_report
from version1.players_app.models import Player



from reportlab.lib import colors
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import mm
from reportlab.platypus import (
    SimpleDocTemplate,
    Paragraph,
    Spacer,
    Table,
    TableStyle,
    PageBreak,
)

from reportlab.lib.pagesizes import A4, landscape




from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


def get_full_name(player):
    """Return the full player name as: First + Second + Surname"""
    if not player:
        return ""
    return f"{player.name} {player.second_name} {player.surname}".strip()


def write_sheet(ws, headers, rows, sum_columns=None):
    """
    Write headers and rows to worksheet with styling and auto-width.
    sum_columns: list of column indices (0-based) to sum at the bottom, e.g., [8,9,10] for Goals, Assists, Pre-Assists
    """
    sum_columns = sum_columns or []

    # Header styling
    header_fill = PatternFill(start_color="0070C0", end_color="0070C0", fill_type="solid")
    header_font = Font(color="FFFFFF", bold=True, size=12, underline="single")  # bold, size 12, underline
    red_border = Border(bottom=Side(style="thin", color="FF0000"))  # red underline

    ws.append(headers)
    for cell in ws[1]:
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        cell.border = red_border

    # Add rows and center them
    for row in rows:
        ws.append(row)
    for row_cells in ws.iter_rows(min_row=2, max_row=ws.max_row, min_col=1, max_col=ws.max_column):
        for cell in row_cells:
            cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)

    # Add total row if sum_columns provided
    if sum_columns and rows:
        totals = ["TOTAL" if i == 0 else "" for i in range(len(headers))]
        for col_idx in sum_columns:
            totals[col_idx] = sum(row[col_idx] if isinstance(row[col_idx], (int, float)) else 0 for row in rows)
        ws.append(totals)
        # Style total row
        for cell in ws[ws.max_row]:
            cell.font = Font(bold=True)
            cell.alignment = Alignment(horizontal="center", vertical="center")

    # Auto-adjust column widths
    for column_cells in ws.columns:
        max_length = 0
        column = column_cells[0].column_letter
        for cell in column_cells:
            if cell.value:
                cell_length = len(str(cell.value))
                if cell_length > max_length:
                    max_length = cell_length
        ws.column_dimensions[column].width = min(max_length + 5, 50)  # max width 50


def download_technical_report(request, team_id):
    team = get_object_or_404(Team, id=team_id)
    season = request.GET.get("season", "")

    start_date_str = request.GET.get("start_date")
    end_date_str = request.GET.get("end_date")
    start_date = parse_date(start_date_str) if start_date_str else None
    end_date = parse_date(end_date_str) if end_date_str else None

    wb = openpyxl.Workbook()
    wb.remove(wb.active)

    def filter_by_date(queryset, date_field='date'):
        if start_date:
            queryset = queryset.filter(**{f"{date_field}__gte": start_date})
        if end_date:
            queryset = queryset.filter(**{f"{date_field}__lte": end_date})
        if season and hasattr(queryset.model, 'season'):
            queryset = queryset.filter(season=season)
        return queryset

    # ================ RESULTS ================
    ws = wb.create_sheet("Results")
    results = filter_by_date(Result.objects.filter(our_team=team))
    
    results_rows = []
    for r in results:
        # Format goal scorers nicely if they are stored as "PlayerName 23, PlayerName 45"
        goal_scorers_text = ""
        if r.goal_scorers:
            scorers_list = [s.strip() for s in r.goal_scorers.split(",")]
            full_names = []
            for s in scorers_list:
                player_name_part = s.rsplit(" ", 1)[0]  # remove minute
                player_obj = Player.objects.filter(name__iexact=player_name_part).first()
                if player_obj:
                    full_names.append(f"{get_full_name(player_obj)} {s.rsplit(' ', 1)[1]}")
                else:
                    full_names.append(s)
            goal_scorers_text = ", ".join(full_names)
        
        results_rows.append([
            r.date,
            r.competition_type,
            r.home_team.name,
            f"{r.home_score}-{r.away_score}",
            r.away_team.name,
            r.result,
            goal_scorers_text
        ])
    
    write_sheet(
        ws,
        ["Date", "Competition", "Home", "Score", "Away", "Result", "Goal Scorers"],
        results_rows
    )

    # ================ MEDICAL ================
    ws = wb.create_sheet("Medical")
    medicals = filter_by_date(Medical.objects.filter(squad=team))
    write_sheet(
        ws,
        ["Player", "Date", "Injury / Illness", "Status", "Comments"],
        [[get_full_name(m.name), m.date, m.injury_or_illness, m.status, m.comments] for m in medicals]
    )

    # ================ TRANSITION ================
    ws = wb.create_sheet("Transition")
    transitions = filter_by_date(Transition.objects.filter(squad=team))
    write_sheet(
        ws,
        ["Player", "Activity", "Played For", "Comments", "Date"],
        [[get_full_name(t.name), t.activity, t.played_for, t.comments, t.date] for t in transitions]
    )

    # ================ SCOUTING ================
    ws = wb.create_sheet("Scouting")
    scouting = filter_by_date(Scouting.objects.filter(squad=team))
    write_sheet(
        ws,
        ["Name", "DOB", "Position", "Agreement", "Former Club", "Comments"],
        [[s.name, s.dob, s.pos, s.agreement, s.former_club, s.comments] for s in scouting]
    )

    # ================ PERFORMANCE ================
    ws = wb.create_sheet("Performance")
    performances = filter_by_date(Performance.objects.filter(squad=team))
    write_sheet(
        ws,
        ["Date", "Activity", "Comments"],
        [[p.date, p.activity, p.comments] for p in performances]
    )

    # ================ INDIVIDUAL ACTION PLAN ================
    ws = wb.create_sheet("Individual Action Plan")
    iaps = filter_by_date(IndividualActionPlan.objects.filter(squad=team))
    write_sheet(
        ws,
        ["Player", "Category", "Responsibility", "Action", "Status", "Follow Up"],
        [[get_full_name(i.name), i.category, i.responsibility, i.action, i.status, i.follow_up] for i in iaps]
    )


def get_full_name(player):
    """Return the full player name as: First + Second + Surname"""
    if not player:
        return ""
    return f"{player.name} {player.second_name} {player.surname}".strip()


def write_sheet(ws, headers, rows, sum_columns=None):
    """
    Write headers and rows to worksheet with styling and auto-width.
    sum_columns: list of column indices (0-based) to sum at the bottom, e.g., [8,9,10] for Goals, Assists, Pre-Assists
    """
    sum_columns = sum_columns or []

    # Header styling
    header_fill = PatternFill(start_color="0070C0", end_color="0070C0", fill_type="solid")
    header_font = Font(color="FFFFFF", bold=True, size=12, underline="single")  # bold, size 12, underline
    red_border = Border(bottom=Side(style="thin", color="FF0000"))  # red underline

    ws.append(headers)
    for cell in ws[1]:
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        cell.border = red_border

    # Add rows and center them
    for row in rows:
        ws.append(row)
    for row_cells in ws.iter_rows(min_row=2, max_row=ws.max_row, min_col=1, max_col=ws.max_column):
        for cell in row_cells:
            cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)

    # Add total row if sum_columns provided
    if sum_columns and rows:
        totals = ["TOTAL" if i == 0 else "" for i in range(len(headers))]
        for col_idx in sum_columns:
            totals[col_idx] = sum(row[col_idx] if isinstance(row[col_idx], (int, float)) else 0 for row in rows)
        ws.append(totals)
        # Style total row
        for cell in ws[ws.max_row]:
            cell.font = Font(bold=True)
            cell.alignment = Alignment(horizontal="center", vertical="center")

    # Auto-adjust column widths
    for column_cells in ws.columns:
        max_length = 0
        column = column_cells[0].column_letter
        for cell in column_cells:
            if cell.value:
                cell_length = len(str(cell.value))
                if cell_length > max_length:
                    max_length = cell_length
        ws.column_dimensions[column].width = min(max_length + 5, 50)  # max width 50


def download_technical_report_excel(request, team_id):
    team = get_object_or_404(Team, id=team_id)
    season = request.GET.get("season", "")

    start_date_str = request.GET.get("start_date")
    end_date_str = request.GET.get("end_date")
    start_date = parse_date(start_date_str) if start_date_str else None
    end_date = parse_date(end_date_str) if end_date_str else None

    wb = openpyxl.Workbook()
    wb.remove(wb.active)

    def filter_by_date(queryset, date_field='date'):
        if start_date:
            queryset = queryset.filter(**{f"{date_field}__gte": start_date})
        if end_date:
            queryset = queryset.filter(**{f"{date_field}__lte": end_date})
        if season and hasattr(queryset.model, 'season'):
            queryset = queryset.filter(season=season)
        return queryset

    # ================ RESULTS ================
    ws = wb.create_sheet("Results")
    results = filter_by_date(Result.objects.filter(our_team=team))
    
    results_rows = []
    for r in results:
        # Format goal scorers nicely if they are stored as "PlayerName 23, PlayerName 45"
        goal_scorers_text = ""
        if r.goal_scorers:
            scorers_list = [s.strip() for s in r.goal_scorers.split(",")]
            full_names = []
            for s in scorers_list:
                player_name_part = s.rsplit(" ", 1)[0]  # remove minute
                player_obj = Player.objects.filter(name__iexact=player_name_part).first()
                if player_obj:
                    full_names.append(f"{get_full_name(player_obj)} {s.rsplit(' ', 1)[1]}")
                else:
                    full_names.append(s)
            goal_scorers_text = ", ".join(full_names)
        
        results_rows.append([
            r.date,
            r.competition_type,
            r.home_team.name,
            f"{r.home_score}-{r.away_score}",
            r.away_team.name,
            r.result,
            goal_scorers_text
        ])
    
    write_sheet(
        ws,
        ["Date", "Competition", "Home", "Score", "Away", "Result", "Goal Scorers"],
        results_rows
    )

    # ================ MEDICAL ================
    ws = wb.create_sheet("Medical")
    medicals = filter_by_date(Medical.objects.filter(squad=team))
    write_sheet(
        ws,
        ["Player", "Date", "Injury / Illness", "Status", "Comments"],
        [[get_full_name(m.name), m.date, m.injury_or_illness, m.status, m.comments] for m in medicals]
    )

    # ================ TRANSITION ================
    ws = wb.create_sheet("Transition")
    transitions = filter_by_date(Transition.objects.filter(squad=team))
    write_sheet(
        ws,
        ["Player", "Activity", "Played For", "Comments", "Date"],
        [[get_full_name(t.name), t.activity, t.played_for, t.comments, t.date] for t in transitions]
    )

    # ================ SCOUTING ================
    ws = wb.create_sheet("Scouting")
    scouting = filter_by_date(Scouting.objects.filter(squad=team))
    write_sheet(
        ws,
        ["Name", "DOB", "Position", "Agreement", "Former Club", "Comments"],
        [[s.name, s.dob, s.pos, s.agreement, s.former_club, s.comments] for s in scouting]
    )

    # ================ PERFORMANCE ================
    ws = wb.create_sheet("Performance")
    performances = filter_by_date(Performance.objects.filter(squad=team))
    write_sheet(
        ws,
        ["Date", "Activity", "Comments"],
        [[p.date, p.activity, p.comments] for p in performances]
    )

    # ================ INDIVIDUAL ACTION PLAN ================
    ws = wb.create_sheet("Individual Action Plan")
    iaps = filter_by_date(IndividualActionPlan.objects.filter(squad=team))
    write_sheet(
        ws,
        ["Player", "Category", "Responsibility", "Action", "Status", "Follow Up"],
        [[get_full_name(i.name), i.category, i.responsibility, i.action, i.status, i.follow_up] for i in iaps]
    )







# ================ STATISTICS ================

    ws = wb.create_sheet("Statistics")
    stats_report = get_statistics_report(filter_type=request.GET.get("filter", "all"),team=team,start_date=start_date,end_date=end_date)

    headers = [
    "NAME","POS","TRAINING TOTAL","U11","U13","U15","U17","U20","FIRST TEAM","NATIONAL TEAM",
    "GAME MINS","APPS","STARTS","SUB IN","SUB OUT","GOALS","ASSISTS","PRE-ASSISTS","NOTE",
    ]

    stats_rows = []

    for r in stats_report:
        player = r.get("player")

        stats_rows.append([
            get_full_name(player),
            str(r.get("position", "")),

            # Training data
            r.get("training_total", 0),r.get("training_u11", 0),r.get("training_u13", 0),r.get("training_u15", 0),r.get("training_u17", 0),r.get("training_u20", 0),
            r.get("training_first", 0),r.get("training_national", 0),

            # Match statistics
            r.get("game_minutes", 0),r.get("appearances", 0),
            r.get("starts", 0),r.get("sub_in", 0),r.get("sub_out", 0),r.get("goals", 0),r.get("assists", 0),r.get("pre_assists", 0),r.get("note", ""),
        ])

        #  totals
        sum_cols = [1,2,3,4,5,6,7,8, 9, 10,11]

    write_sheet(ws,headers,stats_rows,sum_columns=sum_cols)


    # ================ DOWNLOAD ================
    response = HttpResponse(
        content_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    )
    filename = f"Technical_Report_{team.age_group.code}_{season or 'ALL'}.xlsx"
    response["Content-Disposition"] = f'attachment; filename="{filename}"'
    
    wb.save(response)
    return response











def technical_report_page(request, team_id):
    team = get_object_or_404(Team, id=team_id)

    start_date_str = request.GET.get("start_date")
    end_date_str = request.GET.get("end_date")
    season = request.GET.get("season", "")

    start_date = parse_date(start_date_str) if start_date_str else None
    end_date = parse_date(end_date_str) if end_date_str else None

    stats_report = get_statistics_report(
        filter_type=request.GET.get("filter", "all"),
        team=team,
        start_date=start_date,
        end_date=end_date,
    )

    context = {
        "team": team,
        "season": season,
        "start_date": start_date,
        "end_date": end_date,
        "stats_report": stats_report,
    }

    return render(
        request,
        "reports_app/technical_report.html",
        context,
    )



def download_technical_report_pdf(request, team_id):
    team = get_object_or_404(Team, id=team_id)

    start_date_str = request.GET.get("start_date")
    end_date_str = request.GET.get("end_date")
    season = request.GET.get("season", "")

    start_date = parse_date(start_date_str) if start_date_str else None
    end_date = parse_date(end_date_str) if end_date_str else None

    stats_report = get_statistics_report(filter_type=request.GET.get("filter", "all"),team=team,start_date=start_date,end_date=end_date,)
    response = HttpResponse(content_type="application/pdf")

    filename = (
        f"Technical_Report_"
        f"{team.age_group.code}_"
        f"{season or 'ALL'}.pdf"
    )

    response["Content-Disposition"] = (f'attachment; filename="{filename}"')

    document = SimpleDocTemplate(
        response,
        pagesize=landscape(A4),
        rightMargin=10 * mm,
        leftMargin=10 * mm,
        topMargin=10 * mm,
        bottomMargin=10 * mm,
    )

    styles = getSampleStyleSheet()

    story = []

    story.append(
        Paragraph(
            f"TECHNICAL REPORT - {team.name}",
            styles["Title"],
        )
    )

    story.append(
        Paragraph(
            (
                f"Period: "
                f"{start_date or 'All'} "
                f"to "
                f"{end_date or 'All'}"
            ),
            styles["Normal"],
        )
    )

    story.append(Spacer(1, 10))

    table_data = [
    [
    "Player",
    "Pos",
    "Training\nTotal",
    "U11",
    "U13",
    "U15",
    "U17",
    "U20",
    "First\nTeam",
    "National",
    "Game",
    "Apps",
    "Goals",
    "Assists",
    "Pre-Ast",
    ]
    ]

    for row in stats_report:
        player = row.get("player")


        table_data.append([
            get_full_name(player),
            str(row.get("position", "")),

            # Training information
            str(row.get("training_total", 0)),
            str(row.get("training_u11", 0)),
            str(row.get("training_u13", 0)),
            str(row.get("training_u15", 0)),
            str(row.get("training_u17", 0)),
            str(row.get("training_u20", 0)),
            str(row.get("training_first", 0)),
            str(row.get("training_national", 0)),

            # Match information
            str(row.get("game_minutes", 0)),
            str(row.get("appearances", 0)),
            str(row.get("goals", 0)),
            str(row.get("assists", 0)),
            str(row.get("pre_assists", 0)),
        ])


    table = Table(
        table_data,
        repeatRows=1,
    )

    table.setStyle(
        TableStyle([
            (
                "BACKGROUND",
                (0, 0),
                (-1, 0),
                colors.HexColor("#0070C0"),
            ),
            (
                "TEXTCOLOR",
                (0, 0),
                (-1, 0),
                colors.white,
            ),
            (
                "FONTNAME",
                (0, 0),
                (-1, 0),
                "Helvetica-Bold",
            ),
            (
                "ALIGN",
                (0, 0),
                (-1, -1),
                "CENTER",
            ),
            (
                "GRID",
                (0, 0),
                (-1, -1),
                0.5,
                colors.grey,
            ),
            (
                "VALIGN",
                (0, 0),
                (-1, -1),
                "MIDDLE",
            ),
            (
                "FONTSIZE",
                (0, 0),
                (-1, -1),
                8,
            ),
        ])
    )

    story.append(table)

    document.build(story)

    return response


def download_technical_report_ppt(request, team_id):
    team = get_object_or_404(Team, id=team_id)

    start_date_str = request.GET.get("start_date")
    end_date_str = request.GET.get("end_date")
    season = request.GET.get("season", "")

    start_date = parse_date(start_date_str) if start_date_str else None
    end_date = parse_date(end_date_str) if end_date_str else None

    stats_report = get_statistics_report(
        filter_type=request.GET.get("filter", "all"),
        team=team,
        start_date=start_date,
        end_date=end_date,
    )

    presentation = Presentation()

    # -------- TITLE SLIDE --------
    title_slide = presentation.slides.add_slide(
        presentation.slide_layouts[0]
    )

    title_slide.shapes.title.text = "TECHNICAL REPORT"

    title_slide.placeholders[1].text = (
        f"{team.name}\n"
        f"Period: {start_date or 'All'} "
        f"to {end_date or 'All'}"
    )

    # -------- STATISTICS SLIDE --------
    slide = presentation.slides.add_slide(
        presentation.slide_layouts[5]
    )

    title = slide.shapes.title

    if title:
        title.text = "PLAYER STATISTICS"

    rows = min(len(stats_report) + 1, 18)
    columns = 15

    table_shape = slide.shapes.add_table(
        rows,
        columns,
        Inches(0.3),
        Inches(1.2),
        Inches(12.7),
        Inches(5.5),
    )

    table = table_shape.table


    headings = [
    "Player",
    "Pos",
    "Total",
    "U11",
    "U13",
    "U15",
    "U17",
    "U20",
    "First",
    "National",
    "Game",
    "Apps",
    "Goals",
    "Assists",
    "Pre Assists",

    ]

    for column_number, heading in enumerate(headings):
        table.cell(0,column_number).text = heading

    for row_number, row in enumerate(
        stats_report[:17],
        start=1
    ):
        player = row.get("player")


        values = [
            get_full_name(player),
            str(row.get("position", "")),
    
            # Training data
            str(row.get("training_total", 0)),
            str(row.get("training_u11", 0)),
            str(row.get("training_u13", 0)),
            str(row.get("training_u15", 0)),
            str(row.get("training_u17", 0)),
            str(row.get("training_u20", 0)),
            str(row.get("training_first", 0)),
            str(row.get("training_national", 0)),

            # Match data
            str(row.get("game_minutes", 0)),
            str(row.get("appearances", 0)),
            str(row.get("goals", 0)),
            str(row.get("assists", 0)),
            str(row.get("pre_assists", 0)),
            

        ]


        for column_number, value in enumerate(values):
            table.cell(
                row_number,
                column_number
            ).text = value

    response = HttpResponse(
        content_type=(
            "application/vnd.openxmlformats-"
            "officedocument.presentationml.presentation"
        )
    )

    filename = (
        f"Technical_Report_"
        f"{team.age_group.code}_"
        f"{season or 'ALL'}.pptx"
    )

    response["Content-Disposition"] = (
        f'attachment; filename="{filename}"'
    )

    presentation.save(response)

    return response