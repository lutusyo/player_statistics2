from django.http import HttpResponse
from django.shortcuts import get_object_or_404, render
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from pptx import Presentation
from version1.reports_app.models.after_action_review import AfterActionReview
from version1.reports_app.models.weekly_report import WeeklyReport


AAR_FIELDS = [
    ("What Was Planned?", "planned"),
    ("What Actually Happened?", "actual"),
    ("Positives", "positives"),
    ("Negatives", "negatives"),
    ("Reasons for Results", "reasons_for_results"),
    ("Learning Points", "learning_points"),
    ("Player Performance Summary", "player_performance_summary"),
    ("Next Microcycle", "next_microcycle"),
    ("Coach Comments", "coach_comments"),
]

def get_aar(report_id):
    report = get_object_or_404(WeeklyReport, id=report_id)
    aar, _ = AfterActionReview.objects.get_or_create(report=report)
    return report, aar

def after_action_review_detail(request, report_id):
    report, aar = get_aar(report_id)
    return render(request, "reports_app/aar_detail.html", {
    "report": report,
    "sections": [(title, getattr(aar, field) or "No information provided.") for title, field in AAR_FIELDS]
})

def download_aar_pdf(request, report_id):
    report, aar = get_aar(report_id)
    response = HttpResponse(content_type="application/pdf")
    response["Content-Disposition"] = f'attachment; filename="AAR_Week_{report.week}.pdf"'
    styles, story = getSampleStyleSheet(), []
    story += [Paragraph("AFTER ACTION REVIEW", styles["Title"]), Paragraph(f"Team: {report.team} | Week: {report.week}", styles["Heading2"]), Spacer(1, 15)]
    for title, field in AAR_FIELDS:
        story += [Paragraph(title, styles["Heading2"]), Paragraph(getattr(aar, field) or "No information provided.", styles["BodyText"]), Spacer(1, 10)]
    SimpleDocTemplate(response, pagesize=A4).build(story)
    return response

def download_aar_ppt(request, report_id):
    report, aar = get_aar(report_id)
    ppt = Presentation()
    slide = ppt.slides.add_slide(ppt.slide_layouts[0])
    slide.shapes.title.text = "AFTER ACTION REVIEW"
    slide.placeholders[1].text = f"Team: {report.team}\nWeek: {report.week}"
    for title, field in AAR_FIELDS:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = getattr(aar, field) or "No information provided."
    response = HttpResponse(content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")
    response["Content-Disposition"] = f'attachment; filename="AAR_Week_{report.week}.pptx"'
    ppt.save(response)
    return response